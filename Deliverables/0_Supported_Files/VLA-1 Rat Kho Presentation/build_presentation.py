#!/usr/bin/env python3
from __future__ import annotations

import argparse
import io
import json
import re
import zipfile
from pathlib import Path
from typing import Iterable

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Pt
from reportlab.lib import colors as rl_colors
from reportlab.lib.utils import ImageReader, simpleSplit
from reportlab.pdfbase.pdfmetrics import stringWidth
from reportlab.pdfgen import canvas


SCRIPT_DIR = Path(__file__).resolve().parent
REPO_ROOT = SCRIPT_DIR.parents[2]
SPEC_PATH = SCRIPT_DIR / "presentation_spec.json"
EMU_PER_PT = 12700
BULLET = "\u2022"


def emu(value: float) -> int:
    return int(value * EMU_PER_PT)


def parse_hex(color: str) -> tuple[int, int, int]:
    color = color.lstrip("#")
    return tuple(int(color[i : i + 2], 16) for i in (0, 2, 4))


def ppt_color(color: str) -> RGBColor:
    return RGBColor(*parse_hex(color))


def pdf_color(color: str):
    return rl_colors.HexColor(color)


def resolve_path(path_value: str) -> Path:
    path = Path(path_value)
    return path if path.is_absolute() else REPO_ROOT / path


class AssetManager:
    def __init__(self) -> None:
        self._cache: dict[tuple[str, str], Image.Image] = {}

    def _crop_box(self, size: tuple[int, int], crop: dict[str, float] | None):
        width, height = size
        if not crop:
            return (0, 0, width, height)

        left = int(width * crop.get("left", 0.0))
        top = int(height * crop.get("top", 0.0))
        right = int(width * crop.get("right", 1.0))
        bottom = int(height * crop.get("bottom", 1.0))
        right = max(right, left + 1)
        bottom = max(bottom, top + 1)
        return (left, top, right, bottom)

    def image(self, path_value: str, crop: dict[str, float] | None = None) -> Image.Image:
        path = resolve_path(path_value)
        key = (str(path), json.dumps(crop or {}, sort_keys=True))
        if key not in self._cache:
            image = Image.open(path)
            box = self._crop_box(image.size, crop)
            image = image.crop(box)
            if image.mode == "RGBA":
                background = Image.new("RGB", image.size, "white")
                background.paste(image, mask=image.getchannel("A"))
                image = background
            elif image.mode != "RGB":
                image = image.convert("RGB")
            self._cache[key] = image
        return self._cache[key].copy()

    def png_bytes(self, path_value: str, crop: dict[str, float] | None = None) -> io.BytesIO:
        image = self.image(path_value, crop)
        buffer = io.BytesIO()
        image.save(buffer, format="PNG")
        buffer.seek(0)
        return buffer


class PptRenderer:
    def __init__(self, spec: dict, assets: AssetManager) -> None:
        self.spec = spec
        self.assets = assets
        deck = spec["deck"]
        self.width = deck["slide_size_points"]["width"]
        self.height = deck["slide_size_points"]["height"]
        self.theme = spec["theme"]
        self.prs = Presentation()
        self.prs.slide_width = emu(self.width)
        self.prs.slide_height = emu(self.height)
        self.slide = None

    def begin_slide(self) -> None:
        self.slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        background = self.slide.background.fill
        background.solid()
        background.fore_color.rgb = ppt_color(self.theme["background"])

    def save(self, output_path: Path) -> None:
        self.prs.save(str(output_path))

    def panel(
        self,
        x: float,
        y: float,
        width: float,
        height: float,
        *,
        fill: str,
        line: str | None = None,
        line_width: float = 1.0,
        radius: bool = True,
    ):
        shape_type = (
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
            if radius
            else MSO_AUTO_SHAPE_TYPE.RECTANGLE
        )
        shape = self.slide.shapes.add_shape(
            shape_type, emu(x), emu(y), emu(width), emu(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = ppt_color(fill)
        if line:
            shape.line.color.rgb = ppt_color(line)
            shape.line.width = Pt(line_width)
        else:
            shape.line.fill.background()
        return shape

    def text(
        self,
        x: float,
        y: float,
        width: float,
        height: float,
        text: str,
        *,
        font_size: float,
        color: str,
        bold: bool = False,
        align: str = "left",
        valign: str = "top",
        leading: float = 1.15,
        margin: float = 0,
    ):
        box = self.slide.shapes.add_textbox(emu(x), emu(y), emu(width), emu(height))
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_left = emu(margin)
        tf.margin_right = emu(margin)
        tf.margin_top = emu(margin)
        tf.margin_bottom = emu(margin)
        tf.vertical_anchor = {
            "top": MSO_VERTICAL_ANCHOR.TOP,
            "middle": MSO_VERTICAL_ANCHOR.MIDDLE,
            "bottom": MSO_VERTICAL_ANCHOR.BOTTOM,
        }[valign]

        lines = text.split("\n")
        for idx, line in enumerate(lines):
            paragraph = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            paragraph.text = line
            paragraph.alignment = {
                "left": PP_ALIGN.LEFT,
                "center": PP_ALIGN.CENTER,
                "right": PP_ALIGN.RIGHT,
            }[align]
            paragraph.line_spacing = Pt(font_size * leading)
            run = paragraph.runs[0]
            run.font.size = Pt(font_size)
            run.font.name = "Arial"
            run.font.bold = bold
            run.font.color.rgb = ppt_color(color)
        return box

    def bullets(
        self,
        x: float,
        y: float,
        width: float,
        height: float,
        items: Iterable[str],
        *,
        font_size: float,
        color: str,
        leading: float = 1.18,
    ):
        self.text(
            x,
            y,
            width,
            height,
            "\n".join(f"{BULLET} {item}" for item in items),
            font_size=font_size,
            color=color,
            leading=leading,
        )

    def image_box(
        self,
        path_value: str,
        x: float,
        y: float,
        width: float,
        height: float,
        *,
        crop: dict[str, float] | None = None,
        fill: str | None = None,
        line: str | None = None,
        padding: float = 8.0,
    ):
        if fill or line:
            self.panel(
                x,
                y,
                width,
                height,
                fill=fill or self.theme["surface"],
                line=line or self.theme["border"],
                radius=True,
            )

        image = self.assets.image(path_value, crop)
        buffer = self.assets.png_bytes(path_value, crop)
        inner_width = max(width - padding * 2, 1)
        inner_height = max(height - padding * 2, 1)
        image_ratio = image.width / image.height
        box_ratio = inner_width / inner_height

        if image_ratio > box_ratio:
            draw_width = inner_width
            draw_height = draw_width / image_ratio
        else:
            draw_height = inner_height
            draw_width = draw_height * image_ratio

        draw_x = x + (width - draw_width) / 2
        draw_y = y + (height - draw_height) / 2
        self.slide.shapes.add_picture(
            buffer, emu(draw_x), emu(draw_y), emu(draw_width), emu(draw_height)
        )

    def table(
        self,
        x: float,
        y: float,
        width: float,
        height: float,
        headers: list[str],
        rows: list[list[str]],
        *,
        col_ratios: list[float],
        header_fill: str,
        body_fill: str,
        line: str,
        header_font: float = 10.5,
        body_font: float = 9.8,
    ):
        total_ratio = sum(col_ratios)
        col_widths = [width * ratio / total_ratio for ratio in col_ratios]
        row_count = len(rows) + 1
        row_height = height / row_count

        current_x = x
        for idx, header in enumerate(headers):
            cell_width = col_widths[idx]
            self.panel(
                current_x,
                y,
                cell_width,
                row_height,
                fill=header_fill,
                line=line,
                radius=False,
            )
            self.text(
                current_x + 8,
                y + 6,
                cell_width - 16,
                row_height - 12,
                header,
                font_size=header_font,
                color=self.theme["text"],
                bold=True,
                leading=1.05,
            )
            current_x += cell_width

        for row_index, row in enumerate(rows):
            current_x = x
            current_y = y + row_height * (row_index + 1)
            for col_index, cell in enumerate(row):
                cell_width = col_widths[col_index]
                self.panel(
                    current_x,
                    current_y,
                    cell_width,
                    row_height,
                    fill=body_fill,
                    line=line,
                    radius=False,
                )
                self.text(
                    current_x + 8,
                    current_y + 6,
                    cell_width - 16,
                    row_height - 12,
                    cell,
                    font_size=body_font,
                    color=self.theme["text"],
                    leading=1.08,
                )
                current_x += cell_width


class PdfRenderer:
    def __init__(self, spec: dict, assets: AssetManager, output_path: Path) -> None:
        self.spec = spec
        self.assets = assets
        deck = spec["deck"]
        self.width = deck["slide_size_points"]["width"]
        self.height = deck["slide_size_points"]["height"]
        self.theme = spec["theme"]
        self.canvas = canvas.Canvas(str(output_path), pagesize=(self.width, self.height))

    def begin_slide(self) -> None:
        self.canvas.setFillColor(pdf_color(self.theme["background"]))
        self.canvas.rect(0, 0, self.width, self.height, stroke=0, fill=1)

    def save(self) -> None:
        self.canvas.save()

    def _pdf_y(self, y: float, height: float) -> float:
        return self.height - y - height

    def panel(
        self,
        x: float,
        y: float,
        width: float,
        height: float,
        *,
        fill: str,
        line: str | None = None,
        line_width: float = 1.0,
        radius: bool = True,
    ):
        self.canvas.setFillColor(pdf_color(fill))
        if line:
            self.canvas.setStrokeColor(pdf_color(line))
            self.canvas.setLineWidth(line_width)
        else:
            self.canvas.setStrokeColor(pdf_color(fill))
            self.canvas.setLineWidth(0)
        draw_y = self._pdf_y(y, height)
        if radius:
            self.canvas.roundRect(x, draw_y, width, height, 14, stroke=1 if line else 0, fill=1)
        else:
            self.canvas.rect(x, draw_y, width, height, stroke=1 if line else 0, fill=1)

    def text(
        self,
        x: float,
        y: float,
        width: float,
        height: float,
        text: str,
        *,
        font_size: float,
        color: str,
        bold: bool = False,
        align: str = "left",
        leading: float = 1.15,
    ):
        font_name = "Helvetica-Bold" if bold else "Helvetica"
        lines: list[str] = []
        for paragraph in text.split("\n"):
            if paragraph == "":
                lines.append("")
                continue
            lines.extend(simpleSplit(paragraph, font_name, font_size, width))

        line_height = font_size * leading
        top_baseline = self.height - y - font_size
        min_baseline = self.height - (y + height)
        self.canvas.setFillColor(pdf_color(color))
        self.canvas.setFont(font_name, font_size)

        current_y = top_baseline
        for line in lines:
            if current_y < min_baseline:
                break
            if align == "center":
                draw_x = x + (width - stringWidth(line, font_name, font_size)) / 2
            elif align == "right":
                draw_x = x + width - stringWidth(line, font_name, font_size)
            else:
                draw_x = x
            self.canvas.drawString(draw_x, current_y, line)
            current_y -= line_height

    def bullets(
        self,
        x: float,
        y: float,
        width: float,
        height: float,
        items: Iterable[str],
        *,
        font_size: float,
        color: str,
        leading: float = 1.18,
    ):
        font_name = "Helvetica"
        bullet_width = 12
        line_height = font_size * leading
        min_baseline = self.height - (y + height)
        current_y = self.height - y - font_size

        self.canvas.setFillColor(pdf_color(color))
        self.canvas.setFont(font_name, font_size)

        for item in items:
            wrapped = simpleSplit(item, font_name, font_size, width - bullet_width)
            for line_idx, line in enumerate(wrapped):
                if current_y < min_baseline:
                    return
                if line_idx == 0:
                    self.canvas.drawString(x, current_y, BULLET)
                    self.canvas.drawString(x + bullet_width, current_y, line)
                else:
                    self.canvas.drawString(x + bullet_width, current_y, line)
                current_y -= line_height
            current_y -= font_size * 0.15

    def image_box(
        self,
        path_value: str,
        x: float,
        y: float,
        width: float,
        height: float,
        *,
        crop: dict[str, float] | None = None,
        fill: str | None = None,
        line: str | None = None,
        padding: float = 8.0,
    ):
        if fill or line:
            self.panel(
                x,
                y,
                width,
                height,
                fill=fill or self.theme["surface"],
                line=line or self.theme["border"],
                radius=True,
            )

        image = self.assets.image(path_value, crop)
        buffer = self.assets.png_bytes(path_value, crop)
        reader = ImageReader(buffer)
        inner_width = max(width - padding * 2, 1)
        inner_height = max(height - padding * 2, 1)
        image_ratio = image.width / image.height
        box_ratio = inner_width / inner_height

        if image_ratio > box_ratio:
            draw_width = inner_width
            draw_height = draw_width / image_ratio
        else:
            draw_height = inner_height
            draw_width = draw_height * image_ratio

        draw_x = x + (width - draw_width) / 2
        draw_y = y + (height - draw_height) / 2
        self.canvas.drawImage(
            reader,
            draw_x,
            self._pdf_y(draw_y, draw_height),
            width=draw_width,
            height=draw_height,
            preserveAspectRatio=True,
            mask="auto",
        )

    def table(
        self,
        x: float,
        y: float,
        width: float,
        height: float,
        headers: list[str],
        rows: list[list[str]],
        *,
        col_ratios: list[float],
        header_fill: str,
        body_fill: str,
        line: str,
        header_font: float = 10.5,
        body_font: float = 9.8,
    ):
        total_ratio = sum(col_ratios)
        col_widths = [width * ratio / total_ratio for ratio in col_ratios]
        row_count = len(rows) + 1
        row_height = height / row_count

        current_x = x
        for idx, header in enumerate(headers):
            cell_width = col_widths[idx]
            self.panel(
                current_x,
                y,
                cell_width,
                row_height,
                fill=header_fill,
                line=line,
                radius=False,
            )
            self.text(
                current_x + 8,
                y + 6,
                cell_width - 16,
                row_height - 12,
                header,
                font_size=header_font,
                color=self.theme["text"],
                bold=True,
                leading=1.02,
            )
            current_x += cell_width

        for row_index, row in enumerate(rows):
            current_x = x
            current_y = y + row_height * (row_index + 1)
            for col_index, cell in enumerate(row):
                cell_width = col_widths[col_index]
                self.panel(
                    current_x,
                    current_y,
                    cell_width,
                    row_height,
                    fill=body_fill,
                    line=line,
                    radius=False,
                )
                self.text(
                    current_x + 8,
                    current_y + 6,
                    cell_width - 16,
                    row_height - 12,
                    cell,
                    font_size=body_font,
                    color=self.theme["text"],
                    leading=1.05,
                )
                current_x += cell_width

    def end_page(self) -> None:
        self.canvas.showPage()


def render_frame(renderer, title: str, section: str, slide_number: int, total_slides: int, theme: dict):
    renderer.text(42, 26, 620, 28, title, font_size=24, color=theme["text"], bold=True)
    renderer.panel(810, 22, 108, 24, fill=theme["chip_bg"], line=None)
    renderer.text(
        810,
        26,
        108,
        18,
        section.upper(),
        font_size=9.5,
        color=theme["primary_dark"],
        bold=True,
        align="center",
    )
    renderer.panel(42, 72, 876, 1.2, fill=theme["border"], line=None, radius=False)
    renderer.text(42, 516, 280, 12, theme["footer_left"], font_size=8.5, color=theme["muted"])
    renderer.text(
        860,
        516,
        58,
        12,
        f"{slide_number}/{total_slides}",
        font_size=8.5,
        color=theme["muted"],
        align="right",
    )


def render_title(renderer, slide: dict, theme: dict):
    renderer.panel(48, 60, 168, 24, fill=theme["chip_bg"], line=None)
    renderer.text(
        48,
        64,
        168,
        16,
        slide["eyebrow"],
        font_size=9.5,
        color=theme["primary_dark"],
        bold=True,
        align="center",
    )
    renderer.text(48, 108, 420, 68, slide["title"], font_size=30, color=theme["text"], bold=True)
    renderer.text(48, 180, 430, 80, slide["subtitle"], font_size=15, color=theme["muted"], leading=1.22)
    renderer.text(
        48,
        282,
        430,
        120,
        "\n".join(slide["meta_lines"]),
        font_size=13,
        color=theme["text"],
        leading=1.25,
    )
    renderer.image_box(
        slide["image"]["path"],
        536,
        82,
        372,
        286,
        fill=theme["surface"],
        line=theme["border"],
        padding=12,
    )
    renderer.panel(48, 424, 860, 56, fill=theme["surface_alt"], line=theme["border"])
    renderer.text(
        68,
        440,
        820,
        28,
        slide["bottom_note"],
        font_size=12.5,
        color=theme["primary_dark"],
        bold=True,
        align="center",
    )


def render_problem(renderer, slide: dict, theme: dict):
    renderer.bullets(42, 98, 370, 168, slide["bullets"], font_size=14.2, color=theme["text"])
    renderer.panel(438, 98, 220, 182, fill=theme["surface"], line=theme["border"])
    renderer.text(456, 116, 184, 20, slide["left_card"]["title"], font_size=14.5, color=theme["primary_dark"], bold=True)
    renderer.bullets(456, 146, 176, 106, slide["left_card"]["items"], font_size=11.8, color=theme["text"])

    renderer.panel(680, 98, 220, 182, fill=theme["surface"], line=theme["border"])
    renderer.text(698, 116, 184, 20, slide["right_card"]["title"], font_size=14.5, color=theme["red"], bold=True)
    renderer.bullets(698, 146, 176, 106, slide["right_card"]["items"], font_size=11.8, color=theme["text"])

    renderer.panel(42, 312, 858, 140, fill=theme["surface"], line=theme["border"])
    renderer.text(60, 330, 240, 16, "Why this matters operationally", font_size=13, color=theme["primary_dark"], bold=True)
    renderer.text(
        60,
        354,
        824,
        70,
        "A warehouse operator must decide quickly whether the AGV is blocked by harmless material or by a genuine collision risk. The presentation therefore positions the project as decision support for remote assistance, not as autonomous control.",
        font_size=13.2,
        color=theme["text"],
        leading=1.24,
    )
    renderer.panel(60, 430, 824, 1.2, fill=theme["border"], line=None, radius=False)
    renderer.text(60, 440, 824, 18, slide["bottom_note"], font_size=10.8, color=theme["muted"], align="center")


def render_scope(renderer, slide: dict, theme: dict):
    renderer.panel(42, 98, 858, 72, fill=theme["surface_alt"], line=theme["border"])
    renderer.text(60, 114, 120, 18, "Objective", font_size=13, color=theme["primary_dark"], bold=True)
    renderer.text(60, 134, 820, 28, slide["objective"], font_size=13.2, color=theme["text"], leading=1.18)

    renderer.panel(42, 192, 408, 178, fill=theme["surface"], line=theme["border"])
    renderer.text(60, 210, 140, 20, "In Scope", font_size=15, color=theme["green"], bold=True)
    renderer.bullets(60, 240, 360, 110, slide["in_scope"], font_size=12.2, color=theme["text"])

    renderer.panel(492, 192, 408, 178, fill=theme["surface"], line=theme["border"])
    renderer.text(510, 210, 160, 20, "Out of Scope", font_size=15, color=theme["red"], bold=True)
    renderer.bullets(510, 240, 360, 110, slide["out_scope"], font_size=12.2, color=theme["text"])

    card_width = 270
    for idx, contribution in enumerate(slide["contributions"]):
        x = 42 + idx * (card_width + 24)
        renderer.panel(x, 392, card_width, 78, fill=theme["surface"], line=theme["border"])
        renderer.text(
            x + 16,
            414,
            card_width - 32,
            40,
            contribution,
            font_size=13.4,
            color=theme["text"],
            bold=True,
            align="center",
            leading=1.15,
        )


def render_architecture(renderer, slide: dict, theme: dict):
    renderer.panel(42, 98, 260, 340, fill=theme["surface"], line=theme["border"])
    renderer.bullets(60, 120, 220, 220, slide["bullets"], font_size=12.6, color=theme["text"])
    renderer.image_box(
        slide["image"]["path"],
        326,
        98,
        574,
        320,
        fill=theme["surface"],
        line=theme["border"],
        padding=10,
    )
    renderer.text(338, 430, 550, 26, slide["caption"], font_size=10.5, color=theme["muted"], align="center", leading=1.15)


def render_workflow(renderer, slide: dict, theme: dict):
    card_width = 264
    positions = [42, 348, 654]
    for idx, step in enumerate(slide["steps"]):
        x = positions[idx]
        renderer.panel(x, 106, card_width, 282, fill=theme["surface"], line=theme["border"])
        renderer.image_box(
            step["image"]["path"],
            x + 12,
            118,
            card_width - 24,
            170,
            crop=step["image"].get("crop"),
            fill=theme["surface_alt"],
            line=theme["border"],
            padding=4,
        )
        renderer.text(
            x + 16,
            302,
            card_width - 32,
            18,
            step["label"],
            font_size=14,
            color=theme["primary_dark"],
            bold=True,
            align="center",
        )
        renderer.text(
            x + 16,
            328,
            card_width - 32,
            42,
            step["caption"],
            font_size=10.8,
            color=theme["text"],
            align="center",
            leading=1.18,
        )
    renderer.panel(42, 414, 876, 58, fill=theme["surface_alt"], line=theme["border"])
    renderer.text(60, 436, 840, 18, slide["bottom_note"], font_size=12.2, color=theme["primary_dark"], bold=True, align="center")


def render_contract(renderer, slide: dict, theme: dict):
    renderer.image_box(
        slide["image"]["path"],
        42,
        104,
        386,
        278,
        fill=theme["surface"],
        line=theme["border"],
        padding=10,
    )
    renderer.panel(452, 104, 448, 96, fill=theme["surface"], line=theme["border"])
    renderer.text(470, 120, 170, 18, "Model configuration", font_size=14.2, color=theme["primary_dark"], bold=True)
    renderer.bullets(470, 146, 408, 44, slide["config_items"], font_size=10.8, color=theme["text"])

    positions = [(452, 220), (676, 220), (452, 306), (676, 306)]
    for idx, block in enumerate(slide["schema_blocks"]):
        x, y = positions[idx]
        renderer.panel(x, y, 204, 68, fill=theme["surface_alt"], line=theme["border"])
        renderer.text(x + 12, y + 24, 180, 20, block, font_size=12.6, color=theme["text"], bold=True, align="center")

    renderer.panel(42, 406, 858, 64, fill=theme["surface"], line=theme["border"])
    renderer.text(60, 420, 220, 16, "Danger-score mapping", font_size=13, color=theme["primary_dark"], bold=True)
    renderer.panel(300, 418, 250, 28, fill="#DCFCE7", line=None)
    renderer.text(300, 424, 250, 16, slide["score_mapping"][0], font_size=12, color=theme["green"], bold=True, align="center")
    renderer.panel(568, 418, 300, 28, fill="#FEE2E2", line=None)
    renderer.text(568, 424, 300, 16, slide["score_mapping"][1], font_size=12, color=theme["red"], bold=True, align="center")


def render_traceability(renderer, slide: dict, theme: dict):
    renderer.panel(42, 104, 252, 308, fill=theme["surface"], line=theme["border"])
    renderer.bullets(60, 124, 214, 250, slide["bullets"], font_size=12.2, color=theme["text"])
    renderer.image_box(
        slide["main_image"]["path"],
        318,
        104,
        582,
        212,
        crop=slide["main_image"].get("crop"),
        fill=theme["surface"],
        line=theme["border"],
        padding=8,
    )
    renderer.panel(318, 336, 250, 120, fill=theme["surface_alt"], line=theme["border"])
    renderer.text(336, 354, 214, 18, "Traceability value", font_size=13, color=theme["primary_dark"], bold=True)
    renderer.text(
        336,
        380,
        214,
        54,
        "Persisted history snapshots and source IDs make future score-drift analysis auditable instead of anecdotal.",
        font_size=11.2,
        color=theme["text"],
        leading=1.18,
    )
    renderer.image_box(
        slide["inset_image"]["path"],
        586,
        336,
        314,
        120,
        crop=slide["inset_image"].get("crop"),
        fill=theme["surface"],
        line=theme["border"],
        padding=6,
    )


def render_evaluation(renderer, slide: dict, theme: dict):
    renderer.panel(42, 106, 404, 186, fill=theme["surface"], line=theme["border"])
    renderer.text(60, 126, 180, 20, slide["left_stream"]["title"], font_size=15, color=theme["primary_dark"], bold=True)
    renderer.bullets(60, 158, 360, 106, slide["left_stream"]["items"], font_size=12.2, color=theme["text"])

    renderer.panel(494, 106, 406, 186, fill=theme["surface"], line=theme["border"])
    renderer.text(512, 126, 190, 20, slide["right_stream"]["title"], font_size=15, color=theme["primary_dark"], bold=True)
    renderer.bullets(512, 158, 360, 106, slide["right_stream"]["items"], font_size=12.2, color=theme["text"])

    metric_width = 198
    for idx, metric in enumerate(slide["metric_items"]):
        x = 42 + idx * (metric_width + 16)
        renderer.panel(x, 320, metric_width, 64, fill=theme["surface_alt"], line=theme["border"])
        renderer.text(x + 12, 340, metric_width - 24, 30, metric, font_size=11.3, color=theme["text"], bold=True, align="center", leading=1.12)

    renderer.panel(42, 410, 858, 56, fill=theme["surface"], line=theme["border"])
    renderer.text(60, 432, 822, 16, slide["bottom_note"], font_size=12, color=theme["amber"], bold=True, align="center")


def render_audit_results(renderer, slide: dict, theme: dict):
    card_width = 158
    accent_map = {
        "primary": theme["primary"],
        "primary_dark": theme["primary_dark"],
        "amber": theme["amber"],
        "red": theme["red"],
        "green": theme["green"],
    }
    for idx, card in enumerate(slide["cards"]):
        x = 42 + idx * (card_width + 15)
        renderer.panel(x, 112, card_width, 92, fill=theme["surface"], line=theme["border"])
        renderer.panel(x + 12, 124, 134, 6, fill=accent_map[card["accent"]], line=None, radius=False)
        renderer.text(x + 16, 144, 126, 16, card["label"], font_size=11, color=theme["muted"], bold=True, align="center")
        renderer.text(x + 16, 168, 126, 28, card["value"], font_size=25, color=accent_map[card["accent"]], bold=True, align="center")

    renderer.panel(42, 238, 858, 152, fill=theme["surface"], line=theme["border"])
    renderer.text(60, 256, 206, 18, "Interpretation", font_size=13.5, color=theme["primary_dark"], bold=True)
    renderer.bullets(60, 286, 812, 84, slide["bullets"], font_size=12.4, color=theme["text"])

    renderer.panel(42, 414, 858, 48, fill=theme["surface_alt"], line=theme["border"])
    renderer.text(60, 432, 822, 14, slide["source_note"], font_size=10.8, color=theme["muted"], align="center")


def render_repeated_findings(renderer, slide: dict, theme: dict):
    renderer.table(
        42,
        112,
        600,
        214,
        slide["headers"],
        slide["rows"],
        col_ratios=[0.14, 0.11, 0.22, 0.53],
        header_fill=theme["surface_alt"],
        body_fill=theme["surface"],
        line=theme["border"],
    )
    renderer.panel(666, 112, 234, 214, fill=theme["surface"], line=theme["border"])
    renderer.text(684, 132, 198, 18, "Stable groups", font_size=13.5, color=theme["primary_dark"], bold=True)
    renderer.bullets(684, 164, 186, 118, slide["stable_groups"], font_size=11.4, color=theme["text"])

    renderer.panel(42, 354, 858, 102, fill=theme["surface_alt"], line=theme["border"])
    renderer.text(60, 392, 822, 32, slide["bottom_note"], font_size=13, color=theme["primary_dark"], bold=True, align="center", leading=1.15)


def render_limitations(renderer, slide: dict, theme: dict):
    renderer.table(
        42,
        104,
        858,
        320,
        slide["headers"],
        slide["rows"],
        col_ratios=[0.25, 0.28, 0.47],
        header_fill=theme["surface_alt"],
        body_fill=theme["surface"],
        line=theme["border"],
        header_font=10.6,
        body_font=9.3,
    )
    renderer.panel(42, 442, 858, 36, fill=theme["surface_alt"], line=theme["border"])
    renderer.text(60, 454, 822, 14, slide["bottom_note"], font_size=11.2, color=theme["amber"], bold=True, align="center")


def render_conclusion(renderer, slide: dict, theme: dict):
    renderer.panel(42, 112, 394, 220, fill=theme["surface"], line=theme["border"])
    renderer.text(60, 132, 190, 18, "Key takeaways", font_size=14.5, color=theme["primary_dark"], bold=True)
    renderer.bullets(60, 164, 350, 138, slide["takeaways"], font_size=12.4, color=theme["text"])

    renderer.panel(470, 112, 430, 220, fill=theme["surface"], line=theme["border"])
    renderer.text(488, 132, 180, 18, "Future work", font_size=14.5, color=theme["primary_dark"], bold=True)
    roadmap_y = 166
    for idx, item in enumerate(slide["future_work"]):
        renderer.panel(488, roadmap_y + idx * 44, 394, 32, fill=theme["surface_alt"], line=theme["border"])
        renderer.text(500, roadmap_y + 10 + idx * 44, 370, 14, item, font_size=11.4, color=theme["text"], bold=True)

    renderer.panel(42, 366, 858, 86, fill=theme["surface_alt"], line=theme["border"])
    renderer.text(60, 392, 822, 38, slide["bottom_note"], font_size=14, color=theme["primary_dark"], bold=True, align="center", leading=1.18)


def render_references(renderer, slide: dict, theme: dict):
    renderer.panel(42, 112, 858, 314, fill=theme["surface"], line=theme["border"])
    renderer.text(60, 132, 160, 18, "Selected literature", font_size=14.5, color=theme["primary_dark"], bold=True)
    renderer.bullets(60, 164, 804, 180, slide["references"], font_size=12.8, color=theme["text"])
    renderer.panel(42, 446, 858, 32, fill=theme["surface_alt"], line=theme["border"])
    renderer.text(60, 456, 822, 14, slide["note"], font_size=11, color=theme["muted"], align="center")


SLIDE_RENDERERS = {
    "title": render_title,
    "problem": render_problem,
    "scope": render_scope,
    "architecture": render_architecture,
    "workflow": render_workflow,
    "contract": render_contract,
    "traceability": render_traceability,
    "evaluation": render_evaluation,
    "audit_results": render_audit_results,
    "repeated_findings": render_repeated_findings,
    "limitations": render_limitations,
    "conclusion": render_conclusion,
    "references": render_references,
}


def validate_sources(spec: dict) -> None:
    incident_path = resolve_path(
        "Deliverables/0_Supported_Files/VLA-1 Rat Kho Final Report/LaTeX_Source/incident_history_audit.json"
    )
    benchmark_path = resolve_path(
        "Deliverables/0_Supported_Files/VLA-1 Rat Kho Final Report/LaTeX_Source/benchmark_audit.json"
    )
    incident_data = json.loads(incident_path.read_text())
    benchmark_data = json.loads(benchmark_path.read_text())

    for key, expected in spec["validation"]["incident_history_audit"].items():
        actual = incident_data["summary"][key]
        if actual != expected:
            raise ValueError(f"incident_history_audit mismatch for {key}: expected {expected}, found {actual}")

    for key, expected in spec["validation"]["benchmark_audit_overview"].items():
        actual = benchmark_data["overview"][key]
        if actual != expected:
            raise ValueError(f"benchmark_audit mismatch for {key}: expected {expected}, found {actual}")

    required_assets = [
        "Deliverables/3_Functional_Diagram/VLA-1 Rat Kho Architecture Diagram.png",
        "Deliverables/0_Supported_Files/VLA-1 Rat Kho Final Report/LaTeX_Source/figures/system_architecture.png",
        "Deliverables/0_Supported_Files/VLA-1 Rat Kho Final Report/LaTeX_Source/figures/prompt_pipeline.png",
    ]
    for asset in required_assets:
        path = resolve_path(asset)
        if not path.exists():
            raise FileNotFoundError(f"Missing required asset: {path}")


def render_slides(renderer, spec: dict) -> None:
    slides = spec["slides"]
    total = len(slides)
    theme = spec["theme"]
    for idx, slide in enumerate(slides, start=1):
        renderer.begin_slide()
        if slide["type"] != "title":
            render_frame(renderer, slide["title"], slide["section"], idx, total, {**theme, "footer_left": spec["deck"]["footer_left"]})
        render_fn = SLIDE_RENDERERS[slide["type"]]
        render_fn(renderer, slide, theme)
        if slide["type"] == "title":
            renderer.text(48, 516, 280, 12, spec["deck"]["footer_left"], font_size=8.5, color=theme["muted"])
            renderer.text(860, 516, 58, 12, f"{idx}/{total}", font_size=8.5, color=theme["muted"], align="right")
        if isinstance(renderer, PdfRenderer):
            renderer.end_page()


def build_pptx(spec: dict, assets: AssetManager, output_path: Path) -> None:
    renderer = PptRenderer(spec, assets)
    render_slides(renderer, spec)
    renderer.save(output_path)


def build_pdf(spec: dict, assets: AssetManager, output_path: Path) -> None:
    renderer = PdfRenderer(spec, assets, output_path)
    render_slides(renderer, spec)
    renderer.save()


def verify_outputs(output_dir: Path, stem: str) -> dict[str, int]:
    pptx_path = output_dir / f"{stem}.pptx"
    pdf_path = output_dir / f"{stem}.pdf"
    results: dict[str, int] = {}

    if pptx_path.exists():
        with zipfile.ZipFile(pptx_path) as archive:
            slide_entries = [
                name
                for name in archive.namelist()
                if name.startswith("ppt/slides/slide") and name.endswith(".xml")
            ]
        results["pptx_slides"] = len(slide_entries)

    if pdf_path.exists():
        pdf_bytes = pdf_path.read_bytes()
        results["pdf_pages"] = len(re.findall(rb"/Type\s*/Page\b", pdf_bytes))

    return results


def parse_formats(value: str) -> list[str]:
    formats = [item.strip().lower() for item in value.split(",") if item.strip()]
    for item in formats:
        if item not in {"pptx", "pdf"}:
            raise ValueError(f"Unsupported format: {item}")
    return formats


def main() -> None:
    parser = argparse.ArgumentParser(description="Build the AGV Fleet VLA Intelligence presentation deck.")
    parser.add_argument(
        "--formats",
        default="pptx,pdf",
        help="Comma-separated output formats: pptx,pdf",
    )
    parser.add_argument(
        "--output-dir",
        default="../../2_Presentationand_Video",
        help="Output directory for final deliverables, resolved relative to this script.",
    )
    args = parser.parse_args()

    spec = json.loads(SPEC_PATH.read_text())
    validate_sources(spec)

    formats = parse_formats(args.formats)
    output_dir = Path(args.output_dir)
    if not output_dir.is_absolute():
        output_dir = (SCRIPT_DIR / output_dir).resolve()
    output_dir.mkdir(parents=True, exist_ok=True)

    assets = AssetManager()
    stem = spec["deck"]["file_stem"]

    if "pptx" in formats:
        build_pptx(spec, assets, output_dir / f"{stem}.pptx")
    if "pdf" in formats:
        build_pdf(spec, assets, output_dir / f"{stem}.pdf")

    verification = verify_outputs(output_dir, stem)
    print(json.dumps({"output_dir": str(output_dir), "verification": verification}, indent=2))


if __name__ == "__main__":
    main()
